from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # mid navy
HIGHLIGHT = RGBColor(0x0F, 0x3C, 0x96)   # electric blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xE0, 0xE7, 0xFF)   # pale lavender
GREEN     = RGBColor(0x34, 0xD3, 0x99)   # teal‑green
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)   # amber
ORANGE    = RGBColor(0xF9, 0x73, 0x16)   # orange


def fill_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_textbox(slide, text, l, t, w, h, font_size=18, bold=False,
                color: RGBColor = WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def title_slide():
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)
    fill_bg(slide, DARK_BG)

    # accent strip left
    add_rect(slide, 0, 0, 0.35, 7.5, HIGHLIGHT)

    # big heading
    add_textbox(slide, "OpenSplit", 0.7, 1.5, 9, 1.4,
                font_size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # tagline
    add_textbox(slide, "Open‑Source Expense Sharing Platform",
                0.7, 3.1, 10, 0.7, font_size=28, bold=False, color=LIGHT,
                align=PP_ALIGN.LEFT)
    # sub‑tag
    add_textbox(slide,
                "Split bills · Track balances · Settle up — No ads. No paywalls.",
                0.7, 3.95, 10.5, 0.55, font_size=16, color=LIGHT,
                align=PP_ALIGN.LEFT)

    # version / date pill
    add_rect(slide, 0.7, 5.0, 2.8, 0.45, HIGHLIGHT)
    add_textbox(slide, "v1.0  |  May 2026", 0.72, 5.02, 2.8, 0.42,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # decorative circles
    for cx, cy, sz, col in [
        (11.5, 1.0, 2.5, HIGHLIGHT),
        (12.2, 3.5, 1.2, ACCENT),
        (10.8, 5.8, 0.9, HIGHLIGHT),
    ]:
        sh = slide.shapes.add_shape(
            9,  # oval
            Inches(cx), Inches(cy), Inches(sz), Inches(sz)
        )
        sh.fill.solid(); sh.fill.fore_color.rgb = col
        sh.line.fill.background()


def agenda_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Agenda", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    items = [
        ("01", "Project Overview"),
        ("02", "Tech Stack"),
        ("03", "Key Features"),
        ("04", "New Features (v1.1)"),
        ("05", "Application Architecture"),
        ("06", "Data Models"),
        ("07", "API Endpoints"),
        ("08", "Frontend Pages & Routing"),
        ("09", "How to Run the App"),
        ("10", "Summary & Next Steps"),
    ]
    cols = [0.4, 4.6, 8.8]
    row_h = 0.62
    for i, (num, label) in enumerate(items):
        col = cols[i % 3]
        row = 1.3 + (i // 3) * row_h + (i // 3) * 0.08
        add_rect(slide, col, row, 3.9, 0.52, ACCENT)
        add_textbox(slide, num, col + 0.08, row + 0.05, 0.55, 0.42,
                    font_size=15, bold=True, color=YELLOW)
        add_textbox(slide, label, col + 0.7, row + 0.05, 3.1, 0.42,
                    font_size=14, bold=False, color=WHITE)


def overview_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Project Overview", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # left column – description
    add_rect(slide, 0.4, 1.35, 5.8, 5.6, ACCENT)
    add_textbox(slide, "What is OpenSplit?", 0.6, 1.5, 5.4, 0.55,
                font_size=20, bold=True, color=YELLOW)
    desc = (
        "OpenSplit is a free, privacy‑friendly expense\n"
        "sharing web app designed for:\n\n"
        "  •  Students sharing flat expenses\n"
        "  •  Roommates splitting utility bills\n"
        "  •  Friend groups on trips\n"
        "  •  Any shared‑cost scenario\n\n"
        "100% open‑source under the MIT licence.\n"
        "No data sold. No ads. No paywalls."
    )
    add_textbox(slide, desc, 0.6, 2.15, 5.4, 4.5, font_size=15, color=LIGHT)

    # right column – highlights
    add_rect(slide, 6.7, 1.35, 6.2, 5.6, ACCENT)
    add_textbox(slide, "Project Highlights", 6.9, 1.5, 5.8, 0.55,
                font_size=20, bold=True, color=GREEN)
    highlights = [
        ("Target Users",    "Students, roommates, trip groups"),
        ("Open Source",     "MIT licence – fork and extend freely"),
        ("Privacy First",   "Self‑hostable, zero telemetry"),
        ("No Paywalls",     "All features free forever"),
        ("Modern Stack",    "React 18 + Node.js + MongoDB"),
        ("Responsive UI",   "Dark mode & mobile‑friendly"),
    ]
    for j, (title, detail) in enumerate(highlights):
        y = 2.2 + j * 0.72
        add_rect(slide, 6.8, y, 5.9, 0.62, DARK_BG)
        add_textbox(slide, title + ":", 6.95, y + 0.05, 2.0, 0.52,
                    font_size=13, bold=True, color=YELLOW)
        add_textbox(slide, detail, 9.0, y + 0.05, 3.6, 0.52,
                    font_size=13, color=WHITE)


def techstack_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Tech Stack", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cards = [
        ("Frontend",   "React 18 + Vite",          "Component‑based SPA with hot‑reload dev server", GREEN),
        ("UI Library", "MUI (Material Design)",     "Dark/light theming, responsive components, icons", YELLOW),
        ("Backend",    "Node.js + Express",         "RESTful API server with middleware pipeline", ORANGE),
        ("Database",   "MongoDB + Mongoose",        "Document store with ODM schemas & validation", HIGHLIGHT),
        ("Auth",       "JWT (JSON Web Tokens)",     "Stateless 7‑day tokens, bcrypt password hashing", GREEN),
        ("Routing",    "React Router v6",           "Client‑side routing with protected routes", YELLOW),
    ]
    per_row = 3
    card_w, card_h = 3.9, 2.5
    gap_x, gap_y = 0.35, 0.3
    start_x, start_y = 0.35, 1.35
    for i, (layer, tech, desc, col) in enumerate(cards):
        cx = start_x + (i % per_row) * (card_w + gap_x)
        cy = start_y + (i // per_row) * (card_h + gap_y)
        add_rect(slide, cx, cy, card_w, card_h, ACCENT)
        # colour top bar
        add_rect(slide, cx, cy, card_w, 0.12, col)
        add_textbox(slide, layer, cx + 0.12, cy + 0.18, card_w - 0.2, 0.42,
                    font_size=13, bold=True, color=col)
        add_textbox(slide, tech, cx + 0.12, cy + 0.62, card_w - 0.2, 0.45,
                    font_size=17, bold=True, color=WHITE)
        add_textbox(slide, desc, cx + 0.12, cy + 1.15, card_w - 0.25, 1.1,
                    font_size=12, color=LIGHT)


def features_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Key Features", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    features = [
        ("User Auth",            "Register & login with JWT‑secured sessions (7‑day tokens, bcrypt hashing)"),
        ("Group Management",     "Create named expense groups and invite members by email"),
        ("Add Expenses",         "Log shared costs with description, amount, payer and split participants"),
        ("Auto Balances",        "Greedy algorithm computes minimum settlements from net per‑user balances"),
        ("Individual Settlement","Mark each member's share paid separately — no need to settle all at once"),
        ("Expense History",      "Per‑group chronological log with filters (All / Unsettled / Settled)"),
        ("CSV Export",           "Download filtered expense history as a .csv file with one click"),
        ("Micro‑GPT Assistant",  "Built‑in chatbot answers 'who owes what' — no API key, pure algorithm"),
    ]
    for i, (title, detail) in enumerate(features):
        col = 0.35 if i % 2 == 0 else 6.85
        row = 1.35 + (i // 2) * 1.5
        col_accent = GREEN if i % 2 == 0 else YELLOW
        add_rect(slide, col, row, 6.1, 1.3, ACCENT)
        add_rect(slide, col, row, 0.08, 1.3, col_accent)
        add_textbox(slide, title, col + 0.2, row + 0.08, 5.7, 0.42,
                    font_size=15, bold=True, color=col_accent)
        add_textbox(slide, detail, col + 0.2, row + 0.52, 5.7, 0.7,
                    font_size=12, color=LIGHT)


def new_features_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "New Features — v1.1", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cards = [
        (GREEN, "Material Design UI",
         "MUI Components",
         "Full app redesigned with Material Design.\nDark / light theme toggle persisted in localStorage.\nResponsive layouts, proper elevation & theming.\nFixes dark-mode textbox invisibility bug."),
        (YELLOW, "Individual Settlements",
         "Per-Member Mark Paid",
         "Each member's share can be marked paid\nindependently within an expense.\nAuto-closes the expense when all non-payer\nmembers have settled."),
        (ORANGE, "CSV Export",
         "One-Click Download",
         "Export current expense list (any filter:\nAll / Unsettled / Settled) as a .csv file.\nIncludes description, amount, payer, split\nmembers, date and status columns."),
        (HIGHLIGHT, "Micro-GPT Assistant",
         "Built-in Chatbot — No API Key",
         "Floating chat widget on Group Details\nand Settle pages. Keyword-matching algorithm\nqueries live MongoDB data to answer:\n'Who owes what?', 'What's pending?' etc."),
    ]

    positions = [(0.35, 1.35), (6.85, 1.35), (0.35, 4.2), (6.85, 4.2)]
    for (x, y), (col, title, subtitle, body) in zip(positions, cards):
        w, h = 6.1, 2.65
        add_rect(slide, x, y, w, h, ACCENT)
        add_rect(slide, x, y, w, 0.1, col)
        add_rect(slide, x, y + 0.1, 0.38, h - 0.1, col)
        add_textbox(slide, title, x + 0.52, y + 0.18, w - 0.62, 0.42,
                    font_size=16, bold=True, color=col)
        add_textbox(slide, subtitle, x + 0.52, y + 0.6, w - 0.62, 0.32,
                    font_size=12, bold=True, color=YELLOW)
        add_textbox(slide, body, x + 0.52, y + 0.98, w - 0.62, 1.55,
                    font_size=11, color=LIGHT)


def architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Application Architecture", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Three‑tier diagram
    tiers = [
        (0.4,  1.45, 3.7, 5.4, GREEN,     "PRESENTATION TIER",
         "React 18 (Vite)\n\nPages:\n• Login / Register\n• Dashboard\n• Group Details\n• Add Expense\n• Expense History\n• Settlement\n\nLibs: React Router, Axios, Recharts"),
        (4.6,  1.45, 3.7, 5.4, YELLOW,    "APPLICATION TIER",
         "Node.js + Express\n\nMiddleware:\n• CORS\n• JSON body parser\n• JWT auth guard\n\nRoutes:\n• /api/auth\n• /api/groups\n• /api/expenses\n• /api/chat\n\nUtils: balanceCalc, microGPT"),
        (8.8,  1.45, 3.7, 5.4, ORANGE,    "DATA TIER",
         "MongoDB (Mongoose)\n\nModels:\n• User\n• Group\n• Expense\n\nFeatures:\n• Indexed queries\n• Populate refs\n• Timestamps\n• bcrypt pre‑save hook"),
    ]
    for x, y, w, h, col, title, body in tiers:
        add_rect(slide, x, y, w, h, ACCENT)
        add_rect(slide, x, y, w, 0.45, col)
        add_textbox(slide, title, x + 0.1, y + 0.07, w - 0.2, 0.35,
                    font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, body, x + 0.18, y + 0.6, w - 0.3, h - 0.75,
                    font_size=11, color=LIGHT)

    # arrows
    for ax in [4.45, 8.65]:
        add_rect(slide, ax, 3.6, 0.18, 0.08, WHITE)


def data_models_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Data Models", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    models = [
        ("User", GREEN, [
            ("name",      "String",   "required, trim"),
            ("email",     "String",   "unique, lowercase"),
            ("password",  "String",   "bcrypt hashed"),
            ("createdAt", "Date",     "auto timestamp"),
        ]),
        ("Group", YELLOW, [
            ("groupName", "String",     "required, trim"),
            ("members",   "[ObjectId]", "ref → User"),
            ("createdBy", "ObjectId",   "ref → User"),
            ("createdAt", "Date",       "auto timestamp"),
        ]),
        ("Expense", ORANGE, [
            ("description", "String",     "required"),
            ("amount",      "Number",     "min 0.01"),
            ("paidBy",      "ObjectId",   "ref → User"),
            ("splitAmong",  "[ObjectId]", "ref → User"),
            ("groupId",     "ObjectId",   "ref → Group"),
            ("settled",     "Boolean",    "default false"),
            ("date",        "Date",       "default now"),
        ]),
    ]

    col_positions = [0.35, 4.6, 8.85]
    for i, (name, col, fields) in enumerate(models):
        x = col_positions[i]
        w = 4.0
        add_rect(slide, x, 1.35, w, 5.6, ACCENT)
        add_rect(slide, x, 1.35, w, 0.42, col)
        add_textbox(slide, name, x + 0.1, 1.39, w - 0.2, 0.35,
                    font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        # header row
        add_rect(slide, x, 1.77, w, 0.3, DARK_BG)
        add_textbox(slide, "Field", x + 0.1, 1.79, 1.1, 0.26,
                    font_size=10, bold=True, color=col)
        add_textbox(slide, "Type", x + 1.25, 1.79, 1.0, 0.26,
                    font_size=10, bold=True, color=col)
        add_textbox(slide, "Notes", x + 2.3, 1.79, 1.6, 0.26,
                    font_size=10, bold=True, color=col)
        for j, (fname, ftype, fnotes) in enumerate(fields):
            fy = 2.15 + j * 0.62
            bg = ACCENT if j % 2 == 0 else DARK_BG
            add_rect(slide, x + 0.05, fy - 0.04, w - 0.1, 0.58, bg)
            add_textbox(slide, fname,  x + 0.12, fy, 1.05, 0.5, font_size=11, color=WHITE)
            add_textbox(slide, ftype,  x + 1.22, fy, 1.0,  0.5, font_size=11, color=YELLOW)
            add_textbox(slide, fnotes, x + 2.28, fy, 1.65, 0.5, font_size=10, color=LIGHT)


def api_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "API Endpoints", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    endpoints = [
        ("POST",   "/api/auth/register",                  "Register new user",                            GREEN),
        ("POST",   "/api/auth/login",                     "Login – returns JWT token",                    GREEN),
        ("GET",    "/api/groups",                         "List all groups for logged‑in user",           YELLOW),
        ("POST",   "/api/groups",                         "Create a new group",                           YELLOW),
        ("POST",   "/api/groups/:id/members",             "Add member to group by email",                 YELLOW),
        ("GET",    "/api/expenses/:groupId",              "List expenses for a group",                    ORANGE),
        ("POST",   "/api/expenses",                       "Add new expense",                              ORANGE),
        ("PATCH",  "/api/expenses/:id/settle-member",     "Mark one member's share as paid",              ORANGE),
        ("GET",    "/api/expenses/:groupId/balances",     "Computed minimum settlement amounts",          ORANGE),
        ("POST",   "/api/chat/:groupId",                  "Micro-GPT: answer expense/balance questions",  GREEN),
    ]

    # header
    header_y = 1.3
    add_rect(slide, 0.3, header_y, 13.0, 0.35, ACCENT)
    for label, xpos, w in [("Method", 0.35, 1.0), ("Endpoint", 1.45, 5.8), ("Description", 7.35, 5.6)]:
        add_textbox(slide, label, xpos, header_y + 0.04, w, 0.3,
                    font_size=12, bold=True, color=WHITE)

    method_colors = {"GET": GREEN, "POST": YELLOW, "PATCH": ORANGE, "DELETE": RGBColor(0xEF, 0x44, 0x44)}
    for i, (method, path, desc, _) in enumerate(endpoints):
        y = 1.72 + i * 0.52
        row_bg = ACCENT if i % 2 == 0 else DARK_BG
        add_rect(slide, 0.3, y, 13.0, 0.48, row_bg)
        mc = method_colors.get(method, WHITE)
        add_rect(slide, 0.35, y + 0.04, 0.9, 0.38, mc)
        add_textbox(slide, method, 0.35, y + 0.06, 0.9, 0.34,
                    font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, path, 1.4, y + 0.06, 5.85, 0.4, font_size=11, color=LIGHT)
        add_textbox(slide, desc, 7.3, y + 0.06, 5.9, 0.4, font_size=11, color=WHITE)


def frontend_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Frontend Pages & Routing", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    pages = [
        ("/login",                  "Login",          "Email + password form, JWT stored in context",       GREEN,  False),
        ("/register",               "Register",       "Name / email / password sign‑up form",               GREEN,  False),
        ("/dashboard",              "Dashboard",      "List groups; create new group; balance overview",    YELLOW, True),
        ("/groups/:id",             "Group Details",  "Members, balances, expense list for one group",      YELLOW, True),
        ("/groups/:id/add-expense", "Add Expense",    "Form to log a shared cost with split selection",     ORANGE, True),
        ("/groups/:id/history",     "Exp. History",   "Chronological expense log with settled status",      ORANGE, True),
        ("/groups/:id/settle",      "Settlement",     "Computed 'who owes whom' with settle buttons",       GREEN,  True),
    ]

    add_rect(slide, 0.3, 1.3, 12.7, 0.38, ACCENT)
    for label, x, w in [("Route", 0.4, 3.5), ("Page", 3.95, 2.3), ("Description", 6.35, 5.0), ("Auth", 11.45, 1.4)]:
        add_textbox(slide, label, x, 1.33, w, 0.32,
                    font_size=12, bold=True, color=WHITE)

    for i, (route, page, desc, col, protected) in enumerate(pages):
        y = 1.73 + i * 0.72
        bg = ACCENT if i % 2 == 0 else DARK_BG
        add_rect(slide, 0.3, y, 12.7, 0.66, bg)
        add_textbox(slide, route, 0.4, y + 0.08, 3.45, 0.5, font_size=11, color=col)
        add_textbox(slide, page, 3.9, y + 0.08, 2.3, 0.5,
                    font_size=12, bold=True, color=WHITE)
        add_textbox(slide, desc, 6.3, y + 0.08, 5.0, 0.5, font_size=11, color=LIGHT)
        lock = "Protected" if protected else "Public"
        lc = YELLOW if protected else GREEN
        add_textbox(slide, lock, 11.4, y + 0.08, 1.5, 0.5,
                    font_size=11, bold=True, color=lc)


def run_commands_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "How to Run the App", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Prerequisites box
    add_rect(slide, 0.35, 1.35, 12.6, 0.95, ACCENT)
    add_textbox(slide, "Prerequisites", 0.55, 1.4, 4, 0.38,
                font_size=14, bold=True, color=YELLOW)
    add_textbox(slide, "Node.js v18+   •   npm v9+   •   MongoDB (local or Atlas)",
                0.55, 1.78, 12.0, 0.38, font_size=13, color=LIGHT)

    # Step boxes
    steps = [
        ("Step 1 — Clone & Enter the Project",
         "git clone https://github.com/your-org/opensplit.git\ncd opensplit",
         GREEN),
        ("Step 2 — Backend Setup",
         "cd backend\nnpm install\ncp .env.example .env\n# Edit .env → set MONGO_URI and JWT_SECRET\nnpm run dev          # starts on http://localhost:5000",
         YELLOW),
        ("Step 3 — Frontend Setup",
         "cd ../frontend\nnpm install\nnpm run dev          # starts on http://localhost:5173",
         ORANGE),
        ("Step 4 — Open in Browser",
         "Visit → http://localhost:5173\nRegister an account → Create a group → Add expenses!",
         GREEN),
    ]

    col_positions = [0.35, 6.85]
    row_positions = [2.42, 4.82]
    for i, (title, code, col) in enumerate(steps):
        x = col_positions[i % 2]
        y = row_positions[i // 2]
        w, h = 6.05, 2.2
        add_rect(slide, x, y, w, h, ACCENT)
        add_rect(slide, x, y, w, 0.1, col)
        add_textbox(slide, title, x + 0.15, y + 0.16, w - 0.25, 0.42,
                    font_size=13, bold=True, color=col)
        # code block background
        add_rect(slide, x + 0.1, y + 0.65, w - 0.2, h - 0.78, DARK_BG)
        add_textbox(slide, code, x + 0.2, y + 0.7, w - 0.35, h - 0.88,
                    font_size=11, color=WHITE)

    # Optional prod note
    add_rect(slide, 0.35, 7.1, 12.6, 0.28, HIGHLIGHT)
    add_textbox(slide, "Production build:  cd frontend && npm run build   |   Backend: npm start",
                0.5, 7.12, 12.3, 0.24, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


def summary_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 1.15, HIGHLIGHT)
    add_textbox(slide, "Summary & Next Steps", 0.5, 0.15, 12, 0.9,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Summary left
    add_rect(slide, 0.35, 1.35, 6.0, 5.8, ACCENT)
    add_textbox(slide, "What We Built", 0.55, 1.45, 5.6, 0.5,
                font_size=20, bold=True, color=GREEN)
    summary_items = [
        "Full‑stack expense sharing web app",
        "Secure JWT authentication (bcrypt)",
        "CRUD for groups, members & expenses",
        "Greedy minimum‑transaction settlement algorithm",
        "Material Design UI — dark/light theme",
        "Individual per‑member settlement flow",
        "One‑click CSV export of expense history",
        "Built‑in Micro‑GPT chatbot (no API key)",
        "Open‑source (MIT) & self‑hostable",
    ]
    for j, item in enumerate(summary_items):
        add_rect(slide, 0.55, 2.05 + j * 0.55, 0.22, 0.22, GREEN)
        add_textbox(slide, item, 0.88, 2.02 + j * 0.55, 5.3, 0.52,
                    font_size=12, color=LIGHT)

    # Next steps right
    add_rect(slide, 6.9, 1.35, 6.0, 5.8, ACCENT)
    add_textbox(slide, "Potential Next Steps", 7.1, 1.45, 5.6, 0.5,
                font_size=20, bold=True, color=YELLOW)
    next_steps = [
        "Email notifications for settlements",
        "Currency conversion support",
        "Photo receipts / attachment uploads",
        "Recurring expense automation",
        "Mobile app (React Native)",
        "OAuth (Google / GitHub login)",
        "Group chat / activity feed",
        "Deploy to cloud (Render / Vercel)",
    ]
    for j, item in enumerate(next_steps):
        add_rect(slide, 7.1, 2.05 + j * 0.55, 0.22, 0.22, YELLOW)
        add_textbox(slide, item, 7.43, 2.02 + j * 0.55, 5.3, 0.52,
                    font_size=12, color=LIGHT)


def thankyou_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    add_rect(slide, 0, 0, 0.35, 7.5, HIGHLIGHT)
    add_rect(slide, 0, 3.4, 13.33, 0.08, HIGHLIGHT)

    add_textbox(slide, "Thank You!", 0.7, 1.3, 12, 1.3,
                font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "OpenSplit — Making shared expenses simple, transparent & free.",
                0.7, 2.8, 12, 0.6, font_size=20, color=LIGHT, align=PP_ALIGN.CENTER)

    infos = [
        ("Licence",   "MIT — free to use, fork & contribute"),
        ("Frontend",  "http://localhost:5173  (dev)"),
        ("Backend",   "http://localhost:5000  (dev)"),
        ("Stack",     "React 18 · Node.js · MongoDB · JWT"),
    ]
    for j, (label, val) in enumerate(infos):
        x = 1.5 + (j % 2) * 5.5
        y = 4.1 + (j // 2) * 1.0
        add_rect(slide, x, y, 5.0, 0.78, ACCENT)
        add_textbox(slide, label + ":", x + 0.15, y + 0.06, 1.5, 0.3,
                    font_size=12, bold=True, color=YELLOW)
        add_textbox(slide, val, x + 0.15, y + 0.38, 4.7, 0.35,
                    font_size=12, color=LIGHT)

    add_rect(slide, 0.35, 7.1, 12.6, 0.3, HIGHLIGHT)
    add_textbox(slide, "github.com/opensplit  |  MIT Licence  |  Built with Python-pptx",
                0.5, 7.12, 12.3, 0.26, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ── Build all slides ──────────────────────────────────────────────────────────
title_slide()
agenda_slide()
overview_slide()
techstack_slide()
features_slide()
new_features_slide()
architecture_slide()
data_models_slide()
api_slide()
frontend_slide()
run_commands_slide()
summary_slide()
thankyou_slide()

out = "/Users/I587047/Desktop/opensplit/OpenSplit_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
